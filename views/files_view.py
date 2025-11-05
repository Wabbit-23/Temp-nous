import os
import tkinter as tk
import fitz
import sys
import subprocess
from tkinter import ttk, Menu, scrolledtext, messagebox, simpledialog
from pathlib import Path
from typing import Optional
from PIL import Image, ImageTk
from mutagen.mp3 import MP3
from openpyxl import load_workbook
from pptx import Presentation
from docx import Document

from modules.file_manager import FileManager
from modules.file_state import get_starred
from theme.themes import THEMES

class FilesView(ttk.Frame):
    def __init__(self, parent, app_core):
        super().__init__(parent)
        self.app_core = app_core
        self.file_manager = FileManager()
        self.current_directory = os.getcwd()
        self.chat_view = None
        self.current_file = None
        self.file_chip_var = tk.StringVar(value="No file loaded")

        theme = THEMES[self.app_core.current_theme_name]

        self.grid_rowconfigure(1, weight=1)
        self.grid_columnconfigure(0, weight=1)

        self.main_split = ttk.Panedwindow(self, orient=tk.HORIZONTAL)
        self.main_split.pack(fill=tk.BOTH, expand=True)

        self.left_panel = ttk.Frame(self.main_split, width=250, style="Panel.TFrame")
        self.right_panel = ttk.Frame(self.main_split, width=800, style="Panel.TFrame")
        self.main_split.add(self.left_panel, weight=1)
        self.main_split.add(self.right_panel, weight=3)

        self.main_split.bind("<ButtonRelease-1>", self._remember_split_position)
        self.after(200, self._restore_split_position)

        nav = ttk.Frame(self.left_panel, style="Panel.TFrame")
        nav.pack(fill=tk.X, padx=5, pady=5)
        self.button_bar = nav

        self.back_btn = ttk.Button(nav, text="Back", style="FileTop.TButton", command=self.go_back)
        self.back_btn.pack(side=tk.LEFT, padx=(0, 5))

        self.path_entry = ttk.Entry(nav, style="TEntry")
        self.path_entry.pack(side=tk.LEFT, fill=tk.X, expand=True, padx=(0, 5))
        self.path_entry.insert(0, self.current_directory)
        self.path_entry.bind("<Return>", self.jump_to_directory)

        self.search_entry = ttk.Entry(nav, width=20, style="TEntry")
        self.search_entry.pack(side=tk.LEFT, padx=(0, 2))

        self.search_btn = ttk.Button(nav, text="Find", style="FileTop.TButton", command=self.perform_search)
        self.search_btn.pack(side=tk.LEFT, padx=(0, 5))

        self.star_dropdown_btn = ttk.Menubutton(nav, text="Starred", style="TMenubutton")
        self.star_menu = Menu(self.star_dropdown_btn, tearoff=0)
        self.star_dropdown_btn["menu"] = self.star_menu
        self.star_dropdown_btn.bind("<Button-1>", self.refresh_star_menu)
        self.star_dropdown_btn.pack(side=tk.LEFT)

        tree_frame = ttk.Frame(self.left_panel, style="Panel.TFrame")
        tree_frame.pack(fill=tk.BOTH, expand=True)

        self.tree = ttk.Treeview(tree_frame, show="tree")
        self.tree.pack(side=tk.LEFT, fill=tk.BOTH, expand=True)
        self.tree.tag_configure("folder", foreground="skyblue")
        self.tree.tag_configure("file", foreground="white")

        scroll = ttk.Scrollbar(tree_frame, orient="vertical", command=self.tree.yview)
        scroll.pack(side=tk.RIGHT, fill=tk.Y)
        self.tree.configure(yscrollcommand=scroll.set)

        self.tree.bind("<<TreeviewSelect>>", self.on_tree_select)
        self.tree.bind("<Double-1>", self.on_tree_double_click)

        # --- Right-click context menu ---
        self.copied_file_path = None
        self.copied_op = None  # "copy"
        self.context_menu = tk.Menu(self, tearoff=0)
        self.tree.bind("<Button-3>", self.show_context_menu)

        # --- Preview Area ---
        header = ttk.Frame(self.right_panel, style="Card.TFrame", padding=(12, 10))
        header.pack(fill=tk.X, padx=5, pady=(5, 0))
        ttk.Label(header, textvariable=self.file_chip_var, style="StatusChip.TLabel").pack(side=tk.LEFT)

        self.file_text = scrolledtext.ScrolledText(
            self.right_panel, wrap=tk.WORD, relief=tk.FLAT, font=("Segoe UI", 10)
        )
        self.file_text.configure(state='disabled', bg=theme["file_bg"], fg=theme["text"], insertbackground=theme["text"])
        self.file_text.pack(fill=tk.BOTH, expand=True, padx=5, pady=5)

        self.image_preview = ttk.Label(self.right_panel)
        self.image_preview.pack_forget()

        btn_bar = ttk.Frame(self.right_panel)
        btn_bar.pack(fill=tk.X, padx=5, pady=2)
        ttk.Button(btn_bar, text="Save", command=self.save_file).pack(side=tk.LEFT, padx=2)
        ttk.Button(btn_bar, text="New", command=self.new_file).pack(side=tk.LEFT, padx=2)
        ttk.Button(btn_bar, text="Delete", command=self.delete_file).pack(side=tk.LEFT, padx=2)
        ttk.Button(btn_bar, text="Open", command=self.open_file_external).pack(side=tk.LEFT, padx=2)
        ttk.Button(btn_bar, text="Star/Unstar", command=self.toggle_star_selected).pack(side=tk.LEFT, padx=2)

        self.populate_tree()

    # ----- Context Menu -----

    def _build_context_menu(self):
        self.context_menu.delete(0, tk.END)
        self.context_menu.add_command(label="Open File", command=self._context_open_file)
        self.context_menu.add_command(label="Open Preview", command=self._context_open_preview)
        self.context_menu.add_separator()
        self.context_menu.add_command(label="Delete File", command=self._context_delete_file)
        self.context_menu.add_command(label="Star/Unstar", command=self.toggle_star_selected)
        self.context_menu.add_separator()
        self.context_menu.add_command(label="Copy File", command=self._context_copy_file)
        if self.copied_file_path:
            self.context_menu.add_command(label="Paste File", command=self._context_paste_file)

    def show_context_menu(self, event):
        iid = self.tree.identify_row(event.y)
        if iid:
            self.tree.selection_set(iid)
            self._build_context_menu()
            self.context_menu.tk_popup(event.x_root, event.y_root)

    def _context_open_file(self):
        self.open_file_external()

    def _context_open_preview(self):
        sel = self.tree.selection()
        if sel:
            self.load_file(sel[0])

    def _context_delete_file(self):
        self.delete_file()

    def _context_copy_file(self):
        sel = self.tree.selection()
        if sel:
            self.copied_file_path = sel[0]
            self.copied_op = "copy"

    def _context_paste_file(self):
        if not self.copied_file_path:
            return
        sel = self.tree.selection()
        target_dir = sel[0] if sel and os.path.isdir(sel[0]) else self.current_directory
        src = self.copied_file_path
        dst = os.path.join(target_dir, os.path.basename(src))
        import shutil
        try:
            if os.path.isdir(src):
                if os.path.exists(dst):
                    messagebox.showerror("Error", f"Folder already exists: {dst}")
                    return
                shutil.copytree(src, dst)
            else:
                shutil.copy2(src, dst)
            self.populate_tree(target_dir)
            self.copied_file_path = None
            self.copied_op = None
        except Exception as e:
            messagebox.showerror("Error", f"Failed to paste file: {e}")

    # ----- End Context Menu -----

    def _restore_split_position(self):
        ratio = self.app_core.get_file_manager_split() if self.app_core else 0.5
        self.after(0, lambda: self._apply_split_ratio(ratio))

    def _apply_split_ratio(self, ratio):
        if not hasattr(self, 'main_split'):
            return
        self.main_split.update_idletasks()
        total = self.main_split.winfo_width()
        if total <= 0:
            return
        pos = int(total * max(0.1, min(0.9, ratio)))
        try:
            self.main_split.sashpos(0, pos)
        except Exception:
            pass

    def _remember_split_position(self, _event=None):
        if not hasattr(self, 'main_split'):
            return
        total = self.main_split.winfo_width()
        if total <= 0:
            return
        pos = self.main_split.sashpos(0)
        ratio = pos / total if total else 0.5
        if self.app_core:
            self.app_core.set_file_manager_split(ratio)

    def set_chat_view(self, chat_view):
        self.chat_view = chat_view
        self._restore_split_position()
        self._update_status_chip(self.current_file)

    def _update_status_chip(self, path: Optional[str]):
        if not path:
            self.file_chip_var.set("No file loaded")
            return
        p = Path(path)
        display = p.name
        self.file_chip_var.set(f"Loaded - {display}")

    def populate_tree(self, directory=None):
        if directory:
            self.current_directory = directory
        self.path_entry.delete(0, tk.END)
        self.path_entry.insert(0, self.current_directory)
        for iid in self.tree.get_children():
            self.tree.delete(iid)
        try:
            entries = os.listdir(self.current_directory)
        except PermissionError:
            entries = []
        entries.sort(key=str.lower)
        for name in entries:
            full = os.path.join(self.current_directory, name)
            tag = "folder" if os.path.isdir(full) else "file"
            icon = "[DIR]" if os.path.isdir(full) else "[FILE]"
            self.tree.insert("", "end", iid=full, text=f"{icon} {name}", tags=(tag,))

    def go_back(self, event=None):
        parent = os.path.dirname(self.current_directory)
        if os.path.isdir(parent):
            self.populate_tree(parent)

    def jump_to_directory(self, event=None):
        newdir = self.path_entry.get().strip()
        if os.path.isdir(newdir):
            self.populate_tree(newdir)

    def perform_search(self):
        query = self.search_entry.get().strip().lower()
        if not query:
            return

        for iid in self.tree.get_children():
            self.tree.delete(iid)

        matches = []
        for base in self.file_manager.include_paths:
            for root, _, files in os.walk(base):
                for name in files:
                    if query in name.lower():
                        full_path = os.path.join(root, name)
                        matches.append(full_path)
                        if len(matches) >= 100:
                            break
                if len(matches) >= 100:
                    break

        if not matches:
            self.tree.insert("", "end", text="No results found", iid="no_results")
            return

        for path in matches:
            name = os.path.basename(path)
            tag = "folder" if os.path.isdir(path) else "file"
            icon = "[DIR]" if os.path.isdir(path) else "[FILE]"
            self.tree.insert("", "end", iid=path, text=f"{icon} {name}", tags=(tag,))

        first = matches[0]
        self.tree.selection_set(first)
        self.load_file(first)

    def open_file_external(self):
        if not self.current_file:
            return
        try:
            if sys.platform.startswith('win'):
                os.startfile(self.current_file)
            elif sys.platform == 'darwin':
                subprocess.call(('open', self.current_file))
            else:
                subprocess.call(('xdg-open', self.current_file))
        except Exception as e:
            messagebox.showerror("Error", f"Failed to open file: {e}")

    def on_tree_select(self, event=None):
        sel = self.tree.selection()
        if not sel:
            return
        path = sel[0]
        if os.path.isdir(path):
            self.populate_tree(path)
        else:
            self.load_file(path)

    def on_tree_double_click(self, event=None):
        self.on_tree_select()

    def load_file(self, path):
        self.current_file = path
        self._update_status_chip(path)
        ext = Path(path).suffix.lower()
        theme = THEMES[self.app_core.current_theme_name]
        self.file_text.configure(bg=theme["file_bg"], fg=theme["text"], insertbackground=theme["text"])
        self.file_text.pack_forget()
        self.image_preview.pack_forget()
        self.image_preview.configure(background=theme["file_bg"])

        # Image preview
        if ext in {".png", ".jpg", ".jpeg", ".bmp", ".gif", ".tiff", ".webp", ".ico"}:
            try:
                img = Image.open(path)
                img.thumbnail((800, 600))
                self._img_tk = ImageTk.PhotoImage(img)
                self.image_preview.configure(image=self._img_tk, text="", background=theme["file_bg"])
                self.image_preview.pack(fill=tk.BOTH, expand=True, padx=5, pady=5)
            except Exception as e:
                self.image_preview.configure(text=f"[error] Failed to load image: {e}", image="", background=theme["file_bg"])
                self.image_preview.pack(fill=tk.BOTH, expand=True, padx=5, pady=5)

        # PDF preview (first page text)
        elif ext == ".pdf":
            try:
                doc = fitz.open(path)
                text = doc[0].get_text() if doc.page_count > 0 else "[PDF file is empty]"
                self.file_text.configure(state="normal")
                self.file_text.delete("1.0", tk.END)
                self.file_text.insert(tk.END, text.strip() or "[No text found on first page.]")
                self.file_text.configure(state="disabled", bg=theme["file_bg"], fg=theme["text"])
                self.file_text.pack(fill=tk.BOTH, expand=True, padx=5, pady=5)
            except Exception as e:
                self.file_text.configure(state="normal")
                self.file_text.delete("1.0", tk.END)
                self.file_text.insert(tk.END, f"[Error previewing PDF: {e}]")
                self.file_text.configure(state="disabled", bg=theme["file_bg"], fg=theme["text"])
                self.file_text.pack(fill=tk.BOTH, expand=True, padx=5, pady=5)

        # MP3 preview (metadata)
        elif ext == ".mp3":
            try:
                audio = MP3(path)
                duration = audio.info.length
                title = audio.get('TIT2', 'Unknown Title')
                artist = audio.get('TPE1', 'Unknown Artist')
                info = f"MP3 Audio File\nTitle: {title}\nArtist: {artist}\nDuration: {duration:.1f}s"
                self.file_text.configure(state="normal")
                self.file_text.delete("1.0", tk.END)
                self.file_text.insert(tk.END, info)
                self.file_text.configure(state="disabled", bg=theme["file_bg"], fg=theme["text"])
                self.file_text.pack(fill=tk.BOTH, expand=True, padx=5, pady=5)
            except Exception as e:
                self.file_text.configure(state="normal")
                self.file_text.delete("1.0", tk.END)
                self.file_text.insert(tk.END, f"[Error previewing MP3: {e}]")
                self.file_text.configure(state="disabled", bg=theme["file_bg"], fg=theme["text"])
                self.file_text.pack(fill=tk.BOTH, expand=True, padx=5, pady=5)

        # MP4 preview (basic info)
        elif ext == ".mp4":
            try:
                size = os.path.getsize(path) // 1024
                info = f"MP4 Video File\nSize: {size} KB\n(Preview not supported)"
                self.file_text.configure(state="normal")
                self.file_text.delete("1.0", tk.END)
                self.file_text.insert(tk.END, info)
                self.file_text.configure(state="disabled", bg=theme["file_bg"], fg=theme["text"])
                self.file_text.pack(fill=tk.BOTH, expand=True, padx=5, pady=5)
            except Exception as e:
                self.file_text.configure(state="normal")
                self.file_text.delete("1.0", tk.END)
                self.file_text.insert(tk.END, f"[Error previewing MP4: {e}]")
                self.file_text.configure(state="disabled", bg=theme["file_bg"], fg=theme["text"])
                self.file_text.pack(fill=tk.BOTH, expand=True, padx=5, pady=5)

        # Excel (.xlsx) preview (first 10 rows)
        elif ext == ".xlsx":
            try:
                wb = load_workbook(path, read_only=True)
                ws = wb.active
                preview = []
                for i, row in enumerate(ws.iter_rows(values_only=True)):
                    preview.append("\t".join(str(cell) if cell is not None else "" for cell in row))
                    if i >= 9:
                        break
                content = "\n".join(preview)
                self.file_text.configure(state="normal")
                self.file_text.delete("1.0", tk.END)
                self.file_text.insert(tk.END, content)
                self.file_text.configure(state="disabled", bg=theme["file_bg"], fg=theme["text"])
                self.file_text.pack(fill=tk.BOTH, expand=True, padx=5, pady=5)
            except Exception as e:
                self.file_text.configure(state="normal")
                self.file_text.delete("1.0", tk.END)
                self.file_text.insert(tk.END, f"[Error previewing Excel: {e}]")
                self.file_text.configure(state="disabled", bg=theme["file_bg"], fg=theme["text"])
                self.file_text.pack(fill=tk.BOTH, expand=True, padx=5, pady=5)

        # PowerPoint (.pptx) preview (slide titles/text)
        elif ext == ".pptx":
            try:
                prs = Presentation(path)
                slides = []
                for i, slide in enumerate(prs.slides):
                    texts = []
                    for shape in slide.shapes:
                        if hasattr(shape, "text"):
                            texts.append(shape.text)
                    slides.append(f"Slide {i+1}:\n" + "\n".join(texts))
                    if i >= 9:
                        break
                content = "\n\n".join(slides) if slides else "[No slides found]"
                self.file_text.configure(state="normal")
                self.file_text.delete("1.0", tk.END)
                self.file_text.insert(tk.END, content)
                self.file_text.configure(state="disabled", bg=theme["file_bg"], fg=theme["text"])
                self.file_text.pack(fill=tk.BOTH, expand=True, padx=5, pady=5)
            except Exception as e:
                self.file_text.configure(state="normal")
                self.file_text.delete("1.0", tk.END)
                self.file_text.insert(tk.END, f"[Error previewing PowerPoint: {e}]")
                self.file_text.configure(state="disabled", bg=theme["file_bg"], fg=theme["text"])
                self.file_text.pack(fill=tk.BOTH, expand=True, padx=5, pady=5)

        # Text file preview (common formats)
        elif ext in {".txt", ".py", ".md", ".json", ".csv", ".log", ".ini", ".xml", ".html", ".css", ".js", ".ts", ".yaml", ".yml"}:
            try:
                with open(path, "r", encoding="utf-8", errors="ignore") as f:
                    content = f.read()
                self.file_text.configure(state="normal")
                self.file_text.delete("1.0", tk.END)
                self.file_text.insert(tk.END, content)
                self.file_text.configure(state="disabled", bg=theme["file_bg"], fg=theme["text"])
                self.file_text.pack(fill=tk.BOTH, expand=True, padx=5, pady=5)
            except Exception as e:
                self.file_text.configure(state="normal")
                self.file_text.delete("1.0", tk.END)
                self.file_text.insert(tk.END, f"[Error opening file: {e}]")
                self.file_text.configure(state="disabled", bg=theme["file_bg"], fg=theme["text"])
                self.file_text.pack(fill=tk.BOTH, expand=True, padx=5, pady=5)

        elif ext == ".docx":
            try:
                doc = Document(path)
                content = []
                for para in doc.paragraphs:
                    content.append(para.text)
                text = "\n".join(content) or "[No text found in DOCX file.]"
                self.file_text.configure(state="normal")
                self.file_text.delete("1.0", tk.END)
                self.file_text.insert(tk.END, text)
                self.file_text.configure(state="disabled", bg=theme["file_bg"], fg=theme["text"])
                self.file_text.pack(fill=tk.BOTH, expand=True, padx=5, pady=5)
            except Exception as e:
                self.file_text.configure(state="normal")
                self.file_text.delete("1.0", tk.END)
                self.file_text.insert(tk.END, f"[Error previewing DOCX: {e}]")
                self.file_text.configure(state="disabled", bg=theme["file_bg"], fg=theme["text"])
                self.file_text.pack(fill=tk.BOTH, expand=True, padx=5, pady=5)

        else:
            self.file_text.configure(state="normal")
            self.file_text.delete("1.0", tk.END)
            self.file_text.insert(tk.END, f"[Preview not supported for {ext} files.]")
            self.file_text.configure(state="disabled", bg=theme["file_bg"], fg=theme["text"])
            self.file_text.pack(fill=tk.BOTH, expand=True, padx=5, pady=5)

        # AI chat context for text files only
        if self.chat_view:
            try:
                with open(path, "r", encoding="utf-8", errors="ignore") as f:
                    content = f.read()
                self.chat_view.set_file_context(path, content)
            except:
                self.chat_view.set_file_context(path, "")

    def save_file(self):
        if not self.current_file:
            return
        try:
            self.file_text.configure(state="normal")
            content = self.file_text.get("1.0", tk.END)
            with open(self.current_file, "w", encoding="utf-8") as f:
                f.write(content)
            self.file_text.configure(state="disabled")
            if self.app_core:
                self.app_core.show_toast("File saved successfully")
        except Exception as e:
            messagebox.showerror("Error", f"Failed to save file: {e}")

    def new_file(self):
        filename = simpledialog.askstring("New File", "Enter file name:")
        if filename:
            path = os.path.join(self.current_directory, filename)
            with open(path, "w", encoding="utf-8") as f:
                f.write("")
            self.populate_tree()
            self.load_file(path)

    def delete_file(self):
        if not self.current_file:
            return
        import send2trash
        if messagebox.askyesno("Delete File", f"Send {self.current_file} to trash?"):
            send2trash.send2trash(self.current_file)
            self.populate_tree()
            self.file_text.configure(state="normal")
            self.file_text.delete("1.0", tk.END)
            self._update_status_chip(None)
            if self.app_core:
                self.app_core.show_toast("File deleted")
            self.file_text.configure(state="disabled")
            self.image_preview.pack_forget()
            self.current_file = None

    def _open_starred(self, path):
        if os.path.isdir(path):
            self.populate_tree(path)
            return
        # Else, open the file's parent in the tree, then select and preview it
        directory = os.path.dirname(path)
        self.populate_tree(directory)
        if self.tree.exists(path):
            self.tree.selection_set(path)
            self.tree.see(path)
            self.load_file(path)

    def toggle_star_current(self):
        if not self.current_file:
            return
        from modules.file_state import is_starred, add_starred, remove_starred
        if is_starred(self.current_file):
            remove_starred(self.current_file)
        else:
            add_starred(self.current_file)
        self.refresh_star_menu()

    def toggle_star_selected(self):
        sel = self.tree.selection()
        if not sel:
            return
        path = sel[0]
        from modules.file_state import is_starred, add_starred, remove_starred
        if is_starred(path):
            remove_starred(path)
        else:
            add_starred(path)
        self.refresh_star_menu()

    def refresh_star_menu(self, event=None):
        self.star_menu.delete(0, tk.END)
        from modules.file_state import get_starred
        starred = get_starred()
        if not starred:
            self.star_menu.add_command(label="(No starred files/folders)", state="disabled")
            return
        for p in starred:
            name = os.path.basename(p)
            is_folder = os.path.isdir(p)
            icon = "[DIR]" if is_folder else "[FILE]"
            label = f"{icon} {name}"
            self.star_menu.add_command(label=label, command=lambda p=p: self._open_starred(p))
